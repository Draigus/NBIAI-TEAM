from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

ASSETS = r'd:\OneDrive\Claude_code\NBIAI_TEAM\Clients\Couch Heroes\assets'
OUT = r'd:\OneDrive\Claude_code\NBIAI_TEAM\Clients\Couch Heroes\CH_AMA_Deck.pptx'

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
BLUE_PURPLE = RGBColor(0x8B, 0x9F, 0xE8)
TEAL = RGBColor(0x4A, 0xDE, 0xDE)
GOLD = RGBColor(0xF5, 0xC8, 0x42)
BODY_TEXT = RGBColor(0xE8, 0xEA, 0xF0)
MUTED = RGBColor(0x9C, 0xA3, 0xBC)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARKER_BG = RGBColor(0x12, 0x12, 0x20)
CARD_BG = RGBColor(0x24, 0x24, 0x3A)

SW = Inches(10)
SH = Inches(5.625)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

blank_layout = prs.slide_layouts[6]

NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def set_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_full_bleed_image(slide, img_path):
    slide.shapes.add_picture(img_path, 0, 0, SW, SH)


def add_image(slide, img_path, left, top, width=None, height=None):
    kwargs = {}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    slide.shapes.add_picture(img_path, left, top, **kwargs)


def add_dark_overlay(slide, opacity=0.7):
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x18)
    spPr = overlay._element.spPr
    srgb = spPr.find('.//{%s}srgbClr' % NS)
    if srgb is not None:
        alpha_elem = etree.SubElement(srgb, '{%s}alpha' % NS)
        alpha_elem.set('val', str(int((1 - opacity) * 100000)))
    overlay.line.fill.background()


def add_text_box(slide, left, top, width, height, text,
                 font_name='Barlow Condensed', font_size=36,
                 bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_card(slide, left, top, width, height, bg_color=CARD_BG, opacity_pct=85000):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    spPr = shape._element.spPr
    srgb = spPr.find('.//{%s}srgbClr' % NS)
    if srgb is not None:
        alpha_elem = etree.SubElement(srgb, '{%s}alpha' % NS)
        alpha_elem.set('val', str(opacity_pct))
    shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width, color=GOLD):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_full_bleed_image(slide, os.path.join(ASSETS, 'slide8_Picture_12.png'))
add_dark_overlay(slide, 0.6)

add_image(slide, os.path.join(ASSETS, 'slide1_Image_0.png'),
          Inches(2.5), Inches(0.5), width=Inches(5))

add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(1),
             'THE AMA', 'Barlow Condensed', 52, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.6),
             'Your Questions. Our Vision. One Conversation.',
             'Barlow', 22, False, PURPLE, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4),
             'INTERNAL STUDIO COMMUNICATION',
             'Barlow', 10, False, MUTED, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: THE ONE-LINER + WHY THIS GAME
# From Glen's Q6 answer and Q1 closing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
             'YOUR GAME WORLD HOME', 'Barlow Condensed', 44, True, WHITE)

# Glen's Q6 one-liner verbatim
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(1.2),
             'An MMO built around your life, your identity, and the games you love, '
             'not the other way round. Filled with an ever-evolving, ever-alive story '
             'that you can explore.',
             'Barlow', 20, False, GOLD)

# From Q1 closing
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(0.8),
             'The demand is already there. It\u2019s visible in what players play every '
             'day, and in the way the whole industry is quietly shifting.',
             'Barlow', 15, False, BODY_TEXT)

# From Q1
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(7), Inches(1),
             'We want to be part of our players\u2019 lives, not consume them. '
             'That\u2019s the game we\u2019re making, and that\u2019s the home we\u2019re building for them.',
             'Barlow', 15, True, PURPLE)

add_image(slide, os.path.join(ASSETS, 'slide3_Picture_22.png'),
          Inches(7.8), Inches(2.0), height=Inches(3.2))


# ============================================================
# SLIDE 3: WHY THERE IS DEMAND
# From Glen's Q1 answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.4), Inches(1.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             'WHY THIS GAME, WHY NOW', 'Barlow Condensed', 40, True, WHITE)

# Three cards from Q1 answer
card_w = Inches(2.6)
card_h = Inches(3.5)
card_y = Inches(1.5)
gap = Inches(0.35)
start_x = Inches(0.8)

cards_data = [
    ('THE PROBLEM', TEAL,
     'The average WoW player who\u2019d put in more than twenty hours was spending '
     'twenty-eight hours a week in the game. That\u2019s a part-time job. Those players '
     'weren\u2019t always having fun. They were stressing over reputation grinds, gear '
     'checks, raid schedules. Miss a raid night, get kicked from the guild. That\u2019s a '
     'game consuming a person\u2019s life, not being part of it.'),
    ('THE SIGNAL', PURPLE,
     'Fortnite, Roblox, Animal Crossing, Palia, Sky, VRChat, the housing and hangout '
     'side of FFXIV. The numbers are enormous, and the thread running through all of '
     'them is the same. People want somewhere to be. A place to show up, meet their '
     'friends, decorate their corner of it, collect things they care about, and play '
     'around. Progression is expressive, not numerical.'),
    ('OUR ANSWER', GOLD,
     'Couch Heroes is built from the ground up for that. Personalisation and a real '
     'sense of self are a core pillar, not a side feature. Your avatar, your home, '
     'your arcade, your collections. That\u2019s the progression. No gear treadmill. No '
     'forty-hour-a-week commitment. If you\u2019ve got twenty minutes, you drop in and '
     'feel good. If you\u2019ve got an evening, there\u2019s real depth.'),
]

for i, (title, color, desc) in enumerate(cards_data):
    x = start_x + i * (card_w + gap)
    add_card(slide, x, card_y, card_w, card_h)
    add_text_box(slide, x + Inches(0.2), card_y + Inches(0.2),
                 card_w - Inches(0.4), Inches(0.4),
                 title, 'Barlow Condensed', 18, True, color)
    add_text_box(slide, x + Inches(0.2), card_y + Inches(0.7),
                 card_w - Inches(0.4), card_h - Inches(0.9),
                 desc, 'Barlow', 11, False, BODY_TEXT)


# ============================================================
# SLIDE 4: WHAT SETS US APART / USP
# From Glen's Q3 answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.4), Inches(1.5), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             'WHAT SETS US APART', 'Barlow Condensed', 40, True, WHITE)

# From Q3 - the USP breakdown
usp_items = [
    ('AN MMO. FULL STOP.', TEAL,
     'A massively multiplayer online game with a rich and robust role-playing '
     'experience, real combat, real exploration, and a robust social system. '
     'MMO players are who we\u2019re building for.'),
    ('IDENTITY, FRIENDS, COLLECTIONS', PURPLE,
     'You build who you are in our world. Your avatar, your home, your arcade, '
     'your collections. How you dress, how you decorate, the people you hang out '
     'with. Most games treat that as a skin on top of the real game. For us it is '
     'the game.'),
    ('BUILT AROUND YOUR LIFE', GOLD,
     'The game isn\u2019t structured to require grind treadmills or forty-hour weeks. '
     'No chasing a gear score to feel like you belong. No showing up on a Tuesday '
     'night so your guild doesn\u2019t kick you. We\u2019ve designed past those pressures.'),
    ('SKILL OVER STATS', BLUE_PURPLE,
     'What you do matters more than what you\u2019re wearing. Combat, progression, '
     'and mastery come from how you play, not from gear inflation.'),
]

uw = Inches(2.05)
uh = Inches(2.8)
uy = Inches(1.5)
ugap = Inches(0.2)

for i, (title, color, desc) in enumerate(usp_items):
    x = Inches(0.8) + i * (uw + ugap)
    add_card(slide, x, uy, uw, uh)
    add_text_box(slide, x + Inches(0.15), uy + Inches(0.2),
                 uw - Inches(0.3), Inches(0.5),
                 title, 'Barlow Condensed', 12, True, color)
    add_text_box(slide, x + Inches(0.15), uy + Inches(0.75),
                 uw - Inches(0.3), uh - Inches(0.9),
                 desc, 'Barlow', 10.5, False, BODY_TEXT)

# Q3 one-liner at bottom
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.8),
             'Couch Heroes is an MMO you can actually live with, where your identity, '
             'your friends, and your collections carry across the whole gaming world, '
             'and where the rest of the industry plugs in to keep it alive.',
             'Barlow', 12, True, GOLD)


# ============================================================
# SLIDE 5: NOT A WOW KILLER / MARKET POSITION
# From Glen's Q4 and Q7 answers
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.4), Inches(1.5), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             'WE ARE NOT A WOW KILLER', 'Barlow Condensed', 40, True, WHITE)

# From Q4
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(1.0),
             'There\u2019s a massive number of players who loved MMOs but walked away '
             'because the genre asked too much of their lives. They didn\u2019t stop wanting '
             'the experience. They stopped being able to afford the time.',
             'Barlow', 14, False, BODY_TEXT)

add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.6),
             'We\u2019re not asking people to choose us over their existing game. We\u2019re '
             'asking them to come home between sessions of whatever else they play.',
             'Barlow', 14, True, GOLD)

# From Q7 - competitors as validators
add_text_box(slide, Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.4),
             'EVERY COMPETITOR VALIDATES A PIECE OF OUR THESIS',
             'Barlow Condensed', 14, True, TEAL)

comp_items = [
    'FFXIV proves social and housing brilliantly',
    'Roblox proves the platform model works',
    'Fortnite proves the social hub pivot has legs',
    'Palia proves the cosy MMO audience exists',
    'Animal Crossing proves low-pressure expressive gaming sells in the hundreds of millions',
]

for i, item in enumerate(comp_items):
    y = Inches(3.55) + Inches(i * 0.35)
    add_text_box(slide, Inches(1.1), y, Inches(5.5), Inches(0.35),
                 item, 'Barlow', 11, False, BODY_TEXT)
    add_text_box(slide, Inches(0.8), y, Inches(0.3), Inches(0.3),
                 '\u25cf', 'Barlow', 8, False, TEAL)

# From Q7
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.4),
             'None of them combines all of it. That specific combination doesn\u2019t have a direct competitor.',
             'Barlow', 11, True, PURPLE)

add_image(slide, os.path.join(ASSETS, 'slide8_Picture_15.png'),
          Inches(7.0), Inches(0.8), height=Inches(4.2))


# ============================================================
# SLIDE 6: THE PILLARS (TWO LAYERS)
# From Glen's Q9 answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARKER_BG)

add_accent_line(slide, Inches(0.8), Inches(0.4), Inches(1.5), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             'TWO LAYERS OF PILLARS', 'Barlow Condensed', 40, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(1.15), Inches(8), Inches(0.35),
             'The vision pillars tell you why. The foundational pillars tell you what.',
             'Barlow', 13, False, MUTED)

# Vision pillars - top row
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(4), Inches(0.35),
             'VISION PILLARS', 'Barlow Condensed', 16, True, PURPLE)

vision_pillars = [
    'Player identity as the product',
    'Play has real value',
    'Social first, game second',
    'Skill over stats',
    'A living, player-influenced world',
]

vpw = Inches(1.65)
vph = Inches(1.1)
vpy = Inches(2.05)

for i, vp in enumerate(vision_pillars):
    x = Inches(0.55) + i * (vpw + Inches(0.1))
    add_card(slide, x, vpy, vpw, vph)
    add_text_box(slide, x + Inches(0.1), vpy + Inches(0.15),
                 vpw - Inches(0.2), vph - Inches(0.3),
                 vp, 'Barlow', 11, True, PURPLE)

# Foundational pillars - bottom row
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(4), Inches(0.35),
             'FOUNDATIONAL PILLARS', 'Barlow Condensed', 16, True, TEAL)

found_pillars = [
    ('CLASS SYSTEM &\nCUSTOMISATION', TEAL,
     'No locked classes. Your weapons, your build, your approach. '
     'Never stuck with a decision from character creation.'),
    ('PERSISTENT\nDISCOVERY', TEAL,
     'Quests, hidden lore, hidden world. Turn a corner and find '
     'something worth finding, every session.'),
    ('MEANINGFUL\nPLAYER CHOICE', TEAL,
     'How you approach combat, puzzles, your space. '
     'The game responds to your choices.'),
    ('SOCIAL THAT\nRIPPLES OUT', TEAL,
     'Effortless to play with friends. Your progress '
     'and discoveries benefit the people around you.'),
]

fpw = Inches(2.05)
fph = Inches(1.55)
fpy = Inches(3.7)

for i, (title, color, desc) in enumerate(found_pillars):
    x = Inches(0.55) + i * (fpw + Inches(0.15))
    add_card(slide, x, fpy, fpw, fph)
    add_text_box(slide, x + Inches(0.12), fpy + Inches(0.1),
                 fpw - Inches(0.24), Inches(0.45),
                 title, 'Barlow Condensed', 10, True, TEAL)
    add_text_box(slide, x + Inches(0.12), fpy + Inches(0.6),
                 fpw - Inches(0.24), fph - Inches(0.7),
                 desc, 'Barlow', 9.5, False, BODY_TEXT)


# ============================================================
# SLIDE 7: WHO YOU GET TO BE
# From Glen's Q12 answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.4), Inches(1.5), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5.5), Inches(0.7),
             'WHO YOU GET TO BE', 'Barlow Condensed', 40, True, WHITE)

# From Q12
add_text_box(slide, Inches(0.8), Inches(1.25), Inches(5.5), Inches(0.8),
             'In most MMOs, you\u2019re the chosen one. You\u2019re the hero of prophecy. '
             'And that\u2019s great until you\u2019re standing next to 850,000 other chosen ones '
             'all doing the same quest. The fantasy breaks.',
             'Barlow', 13, False, MUTED)

add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.6),
             'We don\u2019t want to hand you one fantasy. We want you to build your own.',
             'Barlow', 16, True, GOLD)

# Three fantasy examples from Q12
fantasies = [
    ('The ultra shadowy assassin who moves through the world like smoke'),
    ('The huge heavy warrior who walks through the front door of every fight'),
    ('The collector who cares more about their shelves than their weapon hand'),
]

for i, f in enumerate(fantasies):
    y = Inches(2.9) + Inches(i * 0.45)
    add_text_box(slide, Inches(1.1), y, Inches(5.5), Inches(0.4),
                 f, 'Barlow', 12, False, BODY_TEXT)
    add_text_box(slide, Inches(0.8), y, Inches(0.3), Inches(0.3),
                 '\u25cf', 'Barlow', 9, False, GOLD)

# From Q12
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1),
             'The game doesn\u2019t tell you who to be. You tell the game. '
             'Your choices land. The corruption you push back stays pushed back. '
             'The people you help remember. The space you build reflects who you\u2019ve become.',
             'Barlow', 12, False, PURPLE)

add_image(slide, os.path.join(ASSETS, 'slide3_Picture_22.png'),
          Inches(7.2), Inches(1.0), height=Inches(3.5))

add_image(slide, os.path.join(ASSETS, 'slide5_Picture_13.png'),
          Inches(8.0), Inches(3.8), height=Inches(1.5))


# ============================================================
# SLIDE 8: WHAT YOU SHOULD FEEL
# From Glen's Q16 answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_full_bleed_image(slide, os.path.join(ASSETS, 'slide8_Picture_12.png'))
add_dark_overlay(slide, 0.6)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             'WHAT YOU SHOULD FEEL', 'Barlow Condensed', 42, True, WHITE)

# Four moods from Q16
moods = [
    ('EXPLORING', TEAL,
     'Wonder. A sense of curiosity pulling you forward. A light in the distance, '
     'a path you haven\u2019t taken, something you want to go and look at just because '
     'it\u2019s there. And agency. Nobody\u2019s telling you where to walk.'),
    ('QUESTING', PURPLE,
     'Properly engaged. You should be telling yourself what you think the next answer '
     'is going to be, building your own theory, only to feel that bit of surprise when '
     'it doesn\u2019t work out exactly like you planned.'),
    ('FIGHTING', GOLD,
     'Fast-paced and fun. Impactful. Grandiose. When you swing, you should feel it '
     'connect. When you pull off something clever with your build, it should feel like '
     'you earned it. Powerful, not stressed.'),
    ('HOME', BLUE_PURPLE,
     'Unhurried. No timer ticking. No daily quest guilt. No anxiety about falling '
     'behind. Just a place you want to be.'),
]

mw = Inches(2.05)
mh = Inches(2.8)
my = Inches(1.5)

for i, (title, color, desc) in enumerate(moods):
    x = Inches(0.8) + i * (mw + Inches(0.2))
    add_card(slide, x, my, mw, mh)
    add_text_box(slide, x + Inches(0.15), my + Inches(0.2),
                 mw - Inches(0.3), Inches(0.4),
                 title, 'Barlow Condensed', 18, True, color)
    add_text_box(slide, x + Inches(0.15), my + Inches(0.7),
                 mw - Inches(0.3), mh - Inches(0.9),
                 desc, 'Barlow', 10.5, False, BODY_TEXT)

# Q16 summary line
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.6),
             'The game should always feel like you chose to be here, not like you have to be. '
             'Every part of the game should earn the player\u2019s attention, not demand it.',
             'Barlow', 12, True, GOLD)


# ============================================================
# SLIDE 9: TRANSITION TO Q&A
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, DARKER_BG)

add_image(slide, os.path.join(ASSETS, 'slide3_Picture_21.png'),
          Inches(7.0), Inches(1.5), height=Inches(3.5))

add_image(slide, os.path.join(ASSETS, 'slide3_Picture_22.png'),
          Inches(7.8), Inches(1.0), height=Inches(3.0))

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(1),
             'YOUR QUESTIONS', 'Barlow Condensed', 52, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5), Inches(0.6),
             "Let\u2019s Talk", 'Barlow', 28, True, PURPLE)

add_text_box(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.2),
             '73 questions across 10 topics.\n'
             'Grouped by theme. Answered one by one.\n'
             'No question too big, none ignored.',
             'Barlow', 14, False, BODY_TEXT)

topics_line = ('Game Vision \u2022 Core Pillars \u2022 Partnerships \u2022 '
               'Company Strategy \u2022 Monetisation \u2022 Roadmap \u2022 '
               'Team \u2022 Culture \u2022 Leadership \u2022 Transparency')
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(8), Inches(0.4),
             topics_line, 'Barlow', 10, False, MUTED)


# ============================================================
# SAVE
# ============================================================
prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Total slides: {len(prs.slides)}')
