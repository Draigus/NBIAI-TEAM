#!/usr/bin/env python3
"""
Generate a project initiative tracking PowerPoint template.
Recreates the Revenue Recognition Rebuild layout with placeholder text.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL    = RGBColor(0x33, 0x33, 0x33)
DARK_GREY   = RGBColor(0x55, 0x55, 0x55)
MID_GREY    = RGBColor(0x88, 0x88, 0x88)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
PALE_GREY   = RGBColor(0xE8, 0xE8, 0xE8)
GREEN_BADGE = RGBColor(0x4C, 0xAF, 0x50)
AMBER_BADGE = RGBColor(0xFF, 0xA0, 0x00)
TEAL        = RGBColor(0x00, 0x7B, 0x83)
BLUE_LINK   = RGBColor(0x1A, 0x73, 0xE8)


# ── Shape helpers ──

def rect(slide, x, y, w, h, fill=None, line_rgb=None, line_w=1):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def oval(slide, x, y, d, fill, line_rgb=None, line_w=1):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def badge(slide, text, x, y, w, h, fill, font_rgb=WHITE, sz=7):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.margin_left = tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = True
    p.font.color.rgb = font_rgb
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return s


# ── Text helpers ──

def tbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def simple_text(slide, text, x, y, w, h, sz=9, bold=False, rgb=BLACK,
                align=PP_ALIGN.LEFT, underline=False):
    tb = tbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = rgb
    p.font.name = "Calibri"
    p.alignment = align
    if underline:
        p.font.underline = True
    return tb


def bullet_list(slide, items, x, y, w, h, sz=7.5):
    tb = tbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_left = Pt(3)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.size = Pt(sz)
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    return tb


def keyed_list(slide, items, x, y, w, h, sz=7):
    tb = tbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_left = Pt(3)
    for i, (prefix, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}  {text}"
        p.font.size = Pt(sz)
        p.font.name = "Calibri"
        p.space_after = Pt(1)
    return tb


# ── Project card ──

def add_card(slide, yb, d):
    """Add one initiative card. yb = y-base, d = card data dict."""

    # Left panel box
    rect(slide, 0.25, yb, 4.15, 2.78, fill=WHITE, line_rgb=LIGHT_GREY, line_w=1.5)

    # Title (bold, underlined)
    simple_text(slide, d["title"], 0.35, yb + 0.06, 3.9, 0.25,
                sz=12, bold=True, underline=True)

    # Description
    simple_text(slide, d["desc"], 0.35, yb + 0.33, 3.9, 0.35, sz=8, rgb=DARK_GREY)

    # LEAD / TEAM badges
    simple_text(slide, "LEAD", 0.35, yb + 0.70, 0.35, 0.18, sz=7, bold=True)
    badge(slide, d["lead"], 0.72, yb + 0.71, 0.9, 0.17,
          fill=PALE_GREY, font_rgb=BLACK, sz=6.5)
    simple_text(slide, "TEAM", 1.72, yb + 0.70, 0.38, 0.18, sz=7, bold=True)
    badge(slide, d["team"], 2.12, yb + 0.71, 1.5, 0.17,
          fill=LIGHT_GREY, font_rgb=BLACK, sz=6.5)

    # Separator line inside panel
    rect(slide, 0.35, yb + 0.95, 3.9, 0.008, fill=LIGHT_GREY)

    # Key Updates
    simple_text(slide, "Key Updates:", 0.35, yb + 0.98, 2.0, 0.18, sz=8, bold=True)
    keyed_list(slide, d["updates"], 0.35, yb + 1.18, 3.9, 1.55, sz=7)

    # ── Status indicators ──
    sx = 4.55

    # STATUS
    simple_text(slide, "STATUS", sx, yb + 0.06, 0.65, 0.18, sz=8, bold=True)
    badge(slide, d["status"], sx + 0.72, yb + 0.06, 0.72, 0.20,
          fill=GREEN_BADGE, font_rgb=WHITE, sz=8)

    # % DONE
    simple_text(slide, "% DONE", sx, yb + 0.32, 0.65, 0.18, sz=8, bold=True)
    rect(slide, sx + 0.72, yb + 0.33, 1.2, 0.17, fill=PALE_GREY)
    pct = d["pct"]
    bar = rect(slide, sx + 0.72, yb + 0.33, max(1.2 * (pct / 100), 0.3), 0.17,
               fill=GREEN_BADGE)
    tf = bar.text_frame
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = f"{pct}%"
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # ETA
    simple_text(slide, "ETA", sx, yb + 0.58, 0.65, 0.18, sz=8, bold=True)
    eb = rect(slide, sx + 0.72, yb + 0.58, 0.85, 0.20,
              fill=WHITE, line_rgb=MID_GREY, line_w=1)
    tf = eb.text_frame
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = d["eta"]
    p.font.size = Pt(7)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # ── Goals ──
    gx = 6.65
    simple_text(slide, "Goals", gx, yb + 0.06, 2.0, 0.20, sz=9, bold=True)
    bullet_list(slide, d["goals"], gx, yb + 0.28, 2.0, 0.80, sz=7.5)

    # ── Dependencies ──
    dx = 8.75
    simple_text(slide, "Dependencies", dx, yb + 0.06, 2.3, 0.20, sz=9, bold=True)
    dy = yb + 0.30
    for dept, sts, desc in d["deps"]:
        simple_text(slide, dept, dx, dy, 0.72, 0.17, sz=7, bold=True)
        c = GREEN_BADGE if sts == "GREEN" else AMBER_BADGE
        badge(slide, sts, dx + 0.72, dy + 0.01, 0.55, 0.15,
              fill=c, font_rgb=WHITE, sz=5.5)
        simple_text(slide, desc, dx + 1.32, dy, 1.0, 0.17, sz=7)
        dy += 0.22

    # ── Risks / Mitigations / Documentation ──
    rx = 11.15
    ry = yb + 0.06

    simple_text(slide, "Risks", rx, ry, 2.0, 0.20, sz=9, bold=True)
    bullet_list(slide, d["risks"], rx, ry + 0.22, 2.0, 0.60, sz=7)

    ry2 = ry + 0.22 + len(d["risks"]) * 0.17 + 0.12
    if d.get("mits"):
        simple_text(slide, "Mitigations", rx, ry2, 2.0, 0.18, sz=9, bold=True)
        bullet_list(slide, d["mits"], rx, ry2 + 0.20, 2.0, 0.40, sz=7)
        ry3 = ry2 + 0.20 + len(d["mits"]) * 0.17 + 0.12
    else:
        ry3 = ry2

    simple_text(slide, "Documentation", rx, ry3, 2.0, 0.18, sz=9, bold=True)
    tb = tbox(slide, rx, ry3 + 0.20, 2.0, 0.40)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    for i, doc in enumerate(d["docs"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {doc}"
        p.font.size = Pt(7)
        p.font.color.rgb = BLUE_LINK
        p.font.name = "Calibri"
        p.font.underline = True
        p.space_after = Pt(2)

    # ── Timeline / Gantt ──
    tx = 4.55
    ty = yb + 1.55
    tw = 6.50

    months = d["months"]
    mw = tw / len(months)

    # Month header labels + vertical grid lines
    for i, m in enumerate(months):
        simple_text(slide, m, tx + i * mw, ty, mw, 0.18, sz=8, bold=True,
                    align=PP_ALIGN.CENTER)
        rect(slide, tx + i * mw, ty + 0.20, 0.005, 0.70, fill=PALE_GREY)

    # Horizontal bar
    bar_y = ty + 0.42
    bar_h = 0.10
    rect(slide, tx + 0.08, bar_y, tw - 0.16, bar_h, fill=TEAL)

    # Milestone dots and labels
    for pct_pos, label, above in d["milestones"]:
        mx = tx + 0.08 + (tw - 0.16) * pct_pos
        dot = 0.11
        oval(slide, mx - dot / 2, bar_y - 0.005, dot, TEAL,
             line_rgb=WHITE, line_w=1.5)
        if above:
            rect(slide, mx - 0.002, bar_y - 0.15, 0.004, 0.15, fill=MID_GREY)
            simple_text(slide, label, mx - 0.55, bar_y - 0.32, 1.10, 0.16,
                        sz=6, rgb=MID_GREY, align=PP_ALIGN.CENTER)
        else:
            rect(slide, mx - 0.002, bar_y + bar_h, 0.004, 0.15, fill=MID_GREY)
            simple_text(slide, label, mx - 0.55, bar_y + bar_h + 0.12, 1.10, 0.16,
                        sz=6, rgb=MID_GREY, align=PP_ALIGN.CENTER)

    # "Project Complete" label at far right
    simple_text(slide, "Project Complete",
                tx + tw - 1.10, bar_y + bar_h + 0.32, 1.20, 0.15,
                sz=6.5, bold=True, rgb=MID_GREY, align=PP_ALIGN.CENTER)


# ── Card data (template placeholders) ──

CARD_1 = dict(
    title="[Initiative 1 Title]",
    desc="[Brief description of this initiative - what needs to happen and why it matters to the business]",
    lead="[Lead Name]",
    team="[Dept 1, Dept 2]",
    status="GREEN",
    pct=40,
    eta="[DD/MM/YY]",
    goals=[
        "[Goal 1 - primary objective]",
        "[Goal 2 - quality or cadence target]",
        "[Goal 3 - delivery commitment]",
    ],
    deps=[
        ("[Dept A]", "GREEN", "[Dependency description]"),
        ("[Dept B]", "GREEN", "[Dependency description]"),
        ("[Dept C]", "GREEN", "[Dependency description]"),
    ],
    risks=["[Risk 1 - data, resource, or timeline risk]"],
    mits=["[Mitigation 1 - action to address risk]"],
    docs=["[Document title or Confluence link]"],
    updates=[
        ("\u2713", "[Completed workstream w/c DD/MM/YY]"),
        ("\u2713", "[Completed workstream w/c DD/MM/YY]"),
        ("\u2022", "[In-progress workstream w/c DD/MM/YY]"),
        ("\u2022", "[Planned workstream w/c DD/MM/YY]"),
        ("\u2022", "[Planned workstream w/c DD/MM/YY]"),
    ],
    months=["Month 1", "Month 2", "Month 3"],
    milestones=[
        (0.05, "[Workstream 1]", False),
        (0.18, "[Workstream 2]", True),
        (0.35, "[Workstream 3]", False),
        (0.48, "[Workstream 4]", True),
        (0.55, "[Workstream 5]", False),
        (0.68, "[Workstream 6]", True),
        (0.78, "[Workstream 7]", False),
        (0.90, "[Workstream 8]", True),
    ],
)

CARD_2 = dict(
    title="[Initiative 2 Title]",
    desc="[Brief description of this initiative - what it delivers and its scope]",
    lead="[Lead Name]",
    team="[Dept 1, Dept 2, Dept 3]",
    status="GREEN",
    pct=65,
    eta="[DD/MM/YY]",
    goals=[
        "[Goal 1 - develop standardised deliverable]",
        "[Goal 2 - replace existing processes]",
    ],
    deps=[
        ("[Body A]", "AMBER", "[Dependency description]"),
        ("[Source]", "AMBER", "[Dependency description]"),
    ],
    risks=[
        "[Risk 1 - alignment or measurement risk]",
        "[Risk 2 - data quality or availability risk]",
        "[Risk 3 - external dependency risk]",
    ],
    mits=[],
    docs=["[Document title]", "[Reference link]"],
    updates=[
        ("\u2713", "[Completed phase w/c DD/MM/YY]"),
        ("\u2022", "[Phase 2 in progress]"),
        ("\u2713", "[Infrastructure setup w/c DD/MM/YY]"),
        ("\u2022", "[Evaluation in progress w/c DD/MM/YY]"),
        ("\u2022", "[Integration ongoing]"),
    ],
    months=["Month 1", "Month 2", "Month 3"],
    milestones=[
        (0.10, "[Phase 1]", False),
        (0.28, "[Phase 2]", True),
        (0.50, "[Phase 3]", False),
        (0.70, "[Phase 4]", True),
        (0.88, "[Phase 5]", False),
    ],
)


# ── Build ──

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    simple_text(slide, "[PROJECT TITLE]", 0.30, 0.10, 12.70, 0.50,
                sz=28, bold=True)

    # Separator bar 1
    rect(slide, 0, 0.68, 13.333, 0.07, fill=CHARCOAL)

    # Card 1
    add_card(slide, 0.83, CARD_1)

    # Separator bar 2
    rect(slide, 0, 3.70, 13.333, 0.12, fill=CHARCOAL)

    # Card 2
    add_card(slide, 3.90, CARD_2)

    return prs


if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "project_initiative_tracker.pptx")
    prs = build()
    prs.save(out_path)
    print(f"Saved: {out_path}")
